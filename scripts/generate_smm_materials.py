from __future__ import annotations

from datetime import date, timedelta
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION_START
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "marketing" / "smm"

PRIMARY = "1E3A5F"
ACCENT = "34D399"
ACCENT_SOFT = "D1FAE5"
TEXT = "1F2937"
MUTED = "6B7280"
LIGHT = "F8FAFC"
BORDER = "D1D5DB"

CLINIC = {
    "name": "Ветеринарная клиника Берлога",
    "tagline": "Старт SMM на 2 недели",
    "city": "Москва",
    "address": "ул. Пилюгина, 6",
    "phone": "+7 (499) 757-34-47",
    "telegram": "t.me/berloga_vet",
    "vk": "vk.com/berloga_vet",
    "site": "berloga.vet",
}

START_DATE = date(2026, 3, 25)
WEEKDAYS_RU = {
    0: "Понедельник",
    1: "Вторник",
    2: "Среда",
    3: "Четверг",
    4: "Пятница",
    5: "Суббота",
    6: "Воскресенье",
}


def hex_rgb(value: str) -> tuple[int, int, int]:
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def doc_rgb(value: str) -> RGBColor:
    return RGBColor(*hex_rgb(value))


def ppt_rgb(value: str) -> PptColor:
    return PptColor(*hex_rgb(value))


def weekday_ru(value: date) -> str:
    return WEEKDAYS_RU[value.weekday()]


def set_cell_bg(cell, color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), color)
    tc_pr.append(shd)


def set_cell_border(cell, color: str = BORDER) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right"):
        tag = f"w:{edge}"
        element = borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), "6")
        element.set(qn("w:space"), "0")
        element.set(qn("w:color"), color)


def style_doc(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Cm(1.5)
    section.bottom_margin = Cm(1.5)
    section.left_margin = Cm(1.8)
    section.right_margin = Cm(1.8)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = doc_rgb(TEXT)

    for style_name, size in (("Title", 22), ("Heading 1", 16), ("Heading 2", 13), ("Heading 3", 11)):
        style = document.styles[style_name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = doc_rgb(PRIMARY)


def add_doc_title(document: Document, title: str, subtitle: str | None = None) -> None:
    heading = document.add_paragraph()
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = heading.add_run(title)
    run.bold = True
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.color.rgb = doc_rgb(PRIMARY)
    if subtitle:
        paragraph = document.add_paragraph()
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        paragraph_run = paragraph.add_run(subtitle)
        paragraph_run.font.name = "Arial"
        paragraph_run.font.size = Pt(10.5)
        paragraph_run.font.color.rgb = doc_rgb(MUTED)


def add_doc_paragraph(document: Document, text: str, bold: bool = False) -> None:
    paragraph = document.add_paragraph()
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "Arial"
    run.font.size = Pt(10.5)
    run.font.color.rgb = doc_rgb(TEXT)
    paragraph.paragraph_format.space_after = Pt(4)


def add_doc_bullets(document: Document, items: list[str]) -> None:
    for item in items:
        paragraph = document.add_paragraph(style="List Bullet")
        run = paragraph.add_run(item)
        run.font.name = "Arial"
        run.font.size = Pt(10.5)
        run.font.color.rgb = doc_rgb(TEXT)


def style_table(table, first_row_fill: str = PRIMARY) -> None:
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            set_cell_border(cell)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(9.5)
                    run.font.color.rgb = doc_rgb(TEXT)
        if row_index == 0:
            for cell in row.cells:
                set_cell_bg(cell, first_row_fill)
                for paragraph in cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
                        run.font.color.rgb = doc_rgb("FFFFFF")


CONTENT_PLAN = [
    {
        "date": START_DATE + timedelta(days=0),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram карусель",
        "topic": "Знакомство с клиникой",
        "goal": "Доверие и узнаваемость",
        "idea": "Кто мы, где находимся, что работаем 24/7, чем отличаемся, какие услуги ключевые.",
        "visual": "Фото клиники, вход, ресепшен, 1 фото команды.",
        "cta": "Написать в сообщения и задать вопрос о питомце.",
    },
    {
        "date": START_DATE + timedelta(days=1),
        "type": "Stories",
        "formats": "Instagram stories, VK stories, Telegram короткие заметки",
        "topic": "Жизнь клиники за день",
        "goal": "Показать живую работу",
        "idea": "3-5 сторис: прием, ресепшен, лаборатория, спокойная атмосфера.",
        "visual": "Вертикальные фото и 1 короткое видео.",
        "cta": "Стикер-опрос: кто у вас дома - кошка или собака?",
    },
    {
        "date": START_DATE + timedelta(days=2),
        "type": "Reels/клип",
        "formats": "Instagram reels, VK клип, Telegram видео",
        "topic": "3 признака, что питомца пора показать врачу сегодня",
        "goal": "Экспертность и сохранения",
        "idea": "Аппетит, вялость, рвота/понос, затрудненное мочеиспускание, одышка.",
        "visual": "Врач в кадре + текстовые титры.",
        "cta": "Сохранить ролик и отправить знакомому владельцу питомца.",
    },
    {
        "date": START_DATE + timedelta(days=3),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram карусель",
        "topic": "Знакомство с врачом",
        "goal": "Личное доверие",
        "idea": "Кратко представить специалиста, опыт, специализацию и подход к пациентам.",
        "visual": "Портрет врача + 2 факта о нем.",
        "cta": "Написать, о каком специалисте рассказать следующим.",
    },
    {
        "date": START_DATE + timedelta(days=4),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram пост",
        "topic": "Отзыв клиента + короткая история пациента",
        "goal": "Социальное доказательство",
        "idea": "Показать отзыв и кратко, как помогли питомцу без сложной медтерминологии.",
        "visual": "Фото питомца после согласования + скрин отзыва.",
        "cta": "Записаться на консультацию через сообщения.",
    },
    {
        "date": START_DATE + timedelta(days=5),
        "type": "Stories",
        "formats": "Instagram stories, VK stories, Telegram заметка",
        "topic": "Мини-викторина",
        "goal": "Вовлеченность",
        "idea": "Вопросы: сколько раз в год вакцинировать, нужен ли чек-ап, можно ли без записи.",
        "visual": "3-4 яркие сторис в фирменных цветах.",
        "cta": "Использовать квиз и реакции.",
    },
    {
        "date": START_DATE + timedelta(days=6),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram карусель",
        "topic": "Услуга недели: УЗИ",
        "goal": "Понятно продать услугу",
        "idea": "Когда нужно УЗИ, что оно показывает, сколько длится, как подготовиться.",
        "visual": "Фото кабинета УЗИ и врача за работой.",
        "cta": "Записаться на УЗИ в директ или сообщения.",
    },
    {
        "date": START_DATE + timedelta(days=7),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram карусель",
        "topic": "Чек-лист после прогулки с собакой",
        "goal": "Полезность и сохранения",
        "idea": "Лапы, клещи, слизистые, аппетит, питье, поведение.",
        "visual": "Список на карточках + фото собаки.",
        "cta": "Сохранить чек-лист.",
    },
    {
        "date": START_DATE + timedelta(days=8),
        "type": "Stories",
        "formats": "Instagram stories, VK stories, Telegram короткий пост",
        "topic": "FAQ по вакцинации",
        "goal": "Снять частые возражения",
        "idea": "Возраст первой прививки, можно ли гулять, нужно ли перед вакциной сдавать анализы.",
        "visual": "Вертикальные карточки с вопросами и ответами.",
        "cta": "Собрать вопросы через стикер.",
    },
    {
        "date": START_DATE + timedelta(days=9),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram пост",
        "topic": "Как работает лаборатория в клинике",
        "goal": "Показать скорость и оборудование",
        "idea": "Расписать, что ОАК за 3 минуты, биохимия за 15 минут и почему это важно.",
        "visual": "Фото лаборатории, пробирок, оборудования.",
        "cta": "Спросить в комментариях, какие анализы чаще всего непонятны владельцам.",
    },
    {
        "date": START_DATE + timedelta(days=10),
        "type": "Reels/клип",
        "formats": "Instagram reels, VK клип, Telegram видео",
        "topic": "Мифы о стерилизации или чистке зубов",
        "goal": "Экспертность и охваты",
        "idea": "1 миф = 1 экран, врач коротко разбирает и дает нормальное объяснение.",
        "visual": "Врач в кадре + подписи крупным шрифтом.",
        "cta": "Сохранить и прислать другу, у кого есть питомец.",
    },
    {
        "date": START_DATE + timedelta(days=11),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram карусель",
        "topic": "Кейс пациента",
        "goal": "Доказать компетентность",
        "idea": "Симптомы, что сделали, какой результат. Без тяжелых подробностей и с понятным языком.",
        "visual": "Фото пациента после разрешения + 3 карточки кейса.",
        "cta": "Написать, если нужна консультация по похожим симптомам.",
    },
    {
        "date": START_DATE + timedelta(days=12),
        "type": "Stories",
        "formats": "Instagram stories, VK stories, Telegram заметка",
        "topic": "Закулисье команды",
        "goal": "Теплый контакт с аудиторией",
        "idea": "Короткие истории о смене, заботе о пациентах, подготовке кабинета.",
        "visual": "Вертикальные фото без перегруза текстом.",
        "cta": "Реакции и быстрый вопрос: какую тему разобрать дальше?",
    },
    {
        "date": START_DATE + timedelta(days=13),
        "type": "Пост",
        "formats": "VK пост, Telegram пост, Instagram пост",
        "topic": "Приглашение задать вопрос ветеринару",
        "goal": "Комментарии и входящие обращения",
        "idea": "Открытый пост: можно оставить вопрос в комментарии или в сообщениях, мы соберем ответы в следующий контент.",
        "visual": "Фото врача + простой заголовок на карточке.",
        "cta": "Оставить вопрос в комментарии или директе.",
    },
]


PRINCIPLES = [
    "Стартуем без перегруза: 4 сильных поста и 3 дня stories в неделю достаточно, если делать это регулярно.",
    "Контент строится по модели 70/20/10: 70% польза и доверие, 20% жизнь клиники и команда, 10% продажи и акции.",
    "Один смысл - три адаптации: основа единая, но под Telegram текст длиннее, под VK добавляем опрос/комментарии, под Instagram делаем визуальные карточки и reels.",
    "Пишем простым человеческим языком: без сложной терминологии, без давления, с заботой и ясным CTA.",
    "Каждый пост должен отвечать хотя бы на один вопрос владельца: когда идти к врачу, к кому записаться, как подготовиться, что делать дома.",
]

PILLARS = [
    "Экспертность: советы врача, разбор симптомов, мифы, чек-листы.",
    "Доверие: врачи, кейсы, отзывы, оборудование, лаборатория, 24/7.",
    "Жизнь клиники: закулисье, смена, атмосфера, пациенты, команда.",
    "Продажи без давления: акции, услуги недели, ответы на частые вопросы, удобная запись.",
]

OPTIMAL_FORMATS = [
    ("Telegram", "3-4 поста в неделю", "Новости, экспертные посты, кейсы, быстрые сервисные сообщения."),
    ("VK", "3-4 поста в неделю + stories", "Посты, клипы, отзывы, обсуждения, прогрев в сообщения."),
    ("Instagram", "3 поста/reels в неделю + stories почти ежедневно", "Визуал, лица врачей, короткие ролики, карусели и сохранения."),
]

REQUIRED_ASSETS = [
    "10-15 хороших фото клиники, врачей, ресепшена, оборудования и пациентов с разрешением.",
    "Список врачей: имя, специализация, опыт, 2 факта о себе, фото.",
    "3-5 отзывов с разрешением на публикацию.",
    "1-2 актуальные акции или услуги, которые важно продвигать прямо сейчас.",
    "Ответственный человек, который 1 раз в неделю согласует темы и фактуру с врачами.",
]

KPI = [
    "Регулярность: минимум 8 основных единиц контента за 14 дней.",
    "Вовлечение: реакции, комментарии, ответы в сторис, сохранения и пересылки.",
    "Лиды: сообщения в Telegram/VK и запросы на запись после контента.",
    "Контент-база: к концу 2 недель должны появиться готовые рубрики, шаблоны и фотоархив.",
]

FORM_TEMPLATES = [
    {
        "title": "Универсальный пост",
        "rows": [
            ("Площадка", "Telegram / VK / Instagram"),
            ("Цель", "Узнаваемость / доверие / запись / прогрев"),
            ("Заголовок", "[Короткий сильный заголовок]"),
            ("Подводка", "[1-2 предложения, почему тема важна]"),
            ("Основной текст", "[Основная мысль, объяснение, польза]"),
            ("CTA", "[Что сделать: написать, сохранить, задать вопрос, записаться]"),
            ("Визуал", "[Какие фото, видео, карточки прикрепить]"),
            ("Текст на обложке", "[3-6 слов для карточки или превью]"),
        ],
    },
    {
        "title": "Экспертный пост врача",
        "rows": [
            ("Тема", "[Например: когда питомцу нужно УЗИ]"),
            ("Врач-эксперт", "[Имя, специализация]"),
            ("Проблема владельца", "[Какой страх или вопрос закрываем]"),
            ("3 ключевых тезиса", "[Тезис 1] / [Тезис 2] / [Тезис 3]"),
            ("Что нельзя делать дома", "[Коротко и понятно]"),
            ("Когда точно нужно ехать в клинику", "[Тревожные признаки]"),
            ("CTA", "[Записаться / задать вопрос]"),
        ],
    },
    {
        "title": "История пациента / кейс",
        "rows": [
            ("Пациент", "[Вид, кличка, возраст]"),
            ("С чем обратились", "[Симптомы без лишних деталей]"),
            ("Что сделали", "[Осмотр, диагностика, лечение]"),
            ("Результат", "[Что стало лучше]"),
            ("Цитата владельца", "[1 фраза]"),
            ("Визуал", "[Фото пациента, врача, 1 карточка кейса]"),
            ("Согласование", "[Есть ли разрешение на публикацию]"),
        ],
    },
    {
        "title": "Пост-знакомство с врачом",
        "rows": [
            ("Имя врача", "[ФИО]"),
            ("Специализация", "[Терапевт / хирург / УЗИ и т.д.]"),
            ("Опыт", "[Сколько лет практики]"),
            ("Чем помогает", "[С какими запросами к нему идут]"),
            ("2 человеческих факта", "[Например: любит кошек, ведет сложные случаи]"),
            ("Фото", "[Портрет + 1 рабочее фото]"),
            ("CTA", "[Задать вопрос врачу / записаться]"),
        ],
    },
    {
        "title": "Продающий пост без давления",
        "rows": [
            ("Услуга или акция", "[Что продвигаем]"),
            ("Для кого", "[Кошки / собаки / щенки / возраст]"),
            ("Почему это важно сейчас", "[Сезонность, выгода, боль владельца]"),
            ("Что входит", "[Коротко по пунктам]"),
            ("Стоимость или оффер", "[Цена / скидка / бонус]"),
            ("Ограничение", "[До какой даты / сколько мест]"),
            ("CTA", "[Написать в сообщения / позвонить]"),
        ],
    },
    {
        "title": "Stories-цепочка на 5 кадров",
        "rows": [
            ("Кадр 1", "[Сильный вопрос или боль владельца]"),
            ("Кадр 2", "[Короткое объяснение]"),
            ("Кадр 3", "[Фото/видео из клиники]"),
            ("Кадр 4", "[Совет или мини-чек-лист]"),
            ("Кадр 5", "[CTA: вопрос, запись, реакция, квиз]"),
            ("Стикеры", "[Опрос / вопрос / шкала]"),
            ("Файлы", "[Какие вертикальные фото или видео нужны]"),
        ],
    },
]


def build_content_plan_docx() -> Path:
    document = Document()
    style_doc(document)
    add_doc_title(
        document,
        "Контент-план на 2 недели",
        f"{CLINIC['name']} • стартовый план на {START_DATE.strftime('%d.%m.%Y')} - {(START_DATE + timedelta(days=13)).strftime('%d.%m.%Y')}",
    )
    add_doc_paragraph(
        document,
        "Этот план собран под старт SMM без перегруза: он помогает быстро показать лицо клиники, экспертность врачей, реальные услуги и мягко приводить людей к записи.",
    )

    document.add_heading("Принципы ведения", level=1)
    add_doc_bullets(document, PRINCIPLES)

    document.add_heading("Оптимальный ритм публикаций", level=1)
    table = document.add_table(rows=1, cols=3)
    headers = ["Площадка", "Частота", "Что публикуем"]
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header
    for platform, frequency, meaning in OPTIMAL_FORMATS:
        row = table.add_row().cells
        row[0].text = platform
        row[1].text = frequency
        row[2].text = meaning
    style_table(table)

    document.add_heading("Детальный план на 14 дней", level=1)
    plan_table = document.add_table(rows=1, cols=8)
    headers = ["Дата", "День", "Тип", "Форматы", "Тема", "Цель", "Визуал", "CTA"]
    for i, header in enumerate(headers):
        plan_table.rows[0].cells[i].text = header
    for item in CONTENT_PLAN:
        row = plan_table.add_row().cells
        row[0].text = item["date"].strftime("%d.%m")
        row[1].text = weekday_ru(item["date"])
        row[2].text = item["type"]
        row[3].text = item["formats"]
        row[4].text = item["topic"]
        row[5].text = f"{item['goal']}. {item['idea']}"
        row[6].text = item["visual"]
        row[7].text = item["cta"]
    style_table(plan_table)

    document.add_heading("Ежедневные stories - база", level=1)
    add_doc_bullets(
        document,
        [
            "1 сторис в начале дня: кто сегодня на смене, что происходит, чем можно помочь.",
            "1 сторис с живым моментом: кабинет, оборудование, прием, врач, спокойный пациент.",
            "1 сторис с вовлечением: опрос, вопрос, квиз, просьба прислать тему для разбора.",
            "1 сервисная сторис по необходимости: свободные окна, напоминание про запись, акция недели.",
        ],
    )

    document.add_heading("Как не выгореть на старте", level=1)
    add_doc_bullets(
        document,
        [
            "Снимать фото и короткие вертикальные видео пакетно 2 раза в неделю.",
            "Один и тот же смысл перерабатывать под 3 площадки, а не придумывать все заново.",
            "Заранее согласовать с врачами 3 экспертные темы на неделю.",
            "Раз в неделю обновлять папку с контентом: врачи, кейсы, отзывы, услуги, акции.",
        ],
    )

    output = OUT_DIR / "Контент-план_2_недели.docx"
    document.save(output)
    return output


def build_templates_docx() -> Path:
    document = Document()
    style_doc(document)
    add_doc_title(
        document,
        "Шаблоны постов и stories",
        f"{CLINIC['name']} • формы для быстрого заполнения",
    )
    add_doc_paragraph(
        document,
        "Используй эти формы как рабочие бланки: копируй нужный шаблон, вставляй текст, фото и комментарии от врача, затем адаптируй под Telegram, VK или Instagram.",
    )

    for template in FORM_TEMPLATES:
        document.add_heading(template["title"], level=1)
        table = document.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "Поле"
        table.rows[0].cells[1].text = "Что вставить"
        for field, value in template["rows"]:
            row = table.add_row().cells
            row[0].text = field
            row[1].text = value
        style_table(table)
        document.add_paragraph()

    document.add_heading("Быстрый weekly-brief перед съемкой", level=1)
    weekly = document.add_table(rows=1, cols=2)
    weekly.rows[0].cells[0].text = "Пункт"
    weekly.rows[0].cells[1].text = "Заполнение"
    for item in [
        "Какие услуги или акции продвигаем на этой неделе",
        "Какие врачи готовы дать комментарий",
        "Какие пациенты/кейсы можно показать",
        "Какие вопросы чаще всего задают владельцы",
        "Какие фото и видео нужно снять отдельно",
    ]:
        row = weekly.add_row().cells
        row[0].text = item
        row[1].text = "[Заполнить]"
    style_table(weekly)

    output = OUT_DIR / "Шаблоны_постов_и_сторис.docx"
    document.save(output)
    return output


def build_fill_forms_md() -> Path:
    content = f"""# Формы для быстрого заполнения

Клиника: {CLINIC['name']}

## 1. Универсальный пост

- Площадка:
- Цель:
- Заголовок:
- Подводка:
- Основной текст:
- CTA:
- Фото/видео:
- Текст на обложке:

## 2. Stories на 5 кадров

- Кадр 1:
- Кадр 2:
- Кадр 3:
- Кадр 4:
- Кадр 5:
- Стикер:
- Какие фото/видео нужны:

## 3. Экспертный пост врача

- Тема:
- Врач:
- Главный вопрос владельца:
- Тезис 1:
- Тезис 2:
- Тезис 3:
- Когда нужно записаться:
- CTA:

## 4. Кейс пациента

- Пациент:
- С чем пришли:
- Что сделали:
- Результат:
- Цитата владельца:
- Фото:
- Разрешение на публикацию:

## 5. Продающий пост

- Услуга/акция:
- Для кого:
- Что входит:
- Почему актуально сейчас:
- Стоимость/оффер:
- Ограничение:
- CTA:
"""
    output = OUT_DIR / "Формы_для_быстрого_заполнения.md"
    output.write_text(content, encoding="utf-8")
    return output


def add_slide_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ppt_rgb(color)


def add_textbox(slide, left, top, width, height, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)
    return box


def add_card(slide, left, top, width, height, title, body, fill=LIGHT, title_color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(fill)
    shape.line.color.rgb = ppt_rgb(BORDER)
    tf = shape.text_frame
    tf.clear()
    title_p = tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = title
    title_run.font.name = "Arial"
    title_run.font.size = PptPt(15)
    title_run.font.bold = True
    title_run.font.color.rgb = ppt_rgb(title_color)
    body_p = tf.add_paragraph()
    body_p.level = 0
    body_run = body_p.add_run()
    body_run.text = body
    body_run.font.name = "Arial"
    body_run.font.size = PptPt(11.5)
    body_run.font.color.rgb = ppt_rgb(TEXT)
    return shape


def add_photo(slide, path: Path, left, top, width, height):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    img_dir = ROOT / "public" / "stock"

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, PRIMARY)
    add_photo(slide, img_dir / "clinic1.jpg", Inches(7.15), Inches(0), Inches(6.18), Inches(7.5))
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.15), Inches(0), Inches(6.18), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = ppt_rgb(PRIMARY)
    overlay.fill.transparency = 0.35
    overlay.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(5.8), Inches(0.7), "SMM-старт для клиники", 18, ACCENT, True)
    add_textbox(slide, Inches(0.8), Inches(1.45), Inches(5.8), Inches(1.3), CLINIC["name"], 28, "FFFFFF", True)
    add_textbox(slide, Inches(0.8), Inches(2.85), Inches(5.8), Inches(1.0), "Контент-план на 2 недели, принципы ведения, оптимальные форматы и готовые шаблоны для работы.", 17, "E5E7EB")
    add_textbox(slide, Inches(0.8), Inches(5.9), Inches(5.5), Inches(1.0), f"{CLINIC['city']} • 24/7 • {CLINIC['telegram']} • {CLINIC['vk']}", 12, ACCENT)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, "FFFFFF")
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.5), "Зачем клинике системный SMM", 24, PRIMARY, True)
    add_card(slide, Inches(0.7), Inches(1.2), Inches(3.9), Inches(2.1), "1. Формировать доверие", "Показывать врачей, подход клиники, оборудование, реальные кейсы и спокойную атмосферу.")
    add_card(slide, Inches(4.75), Inches(1.2), Inches(3.9), Inches(2.1), "2. Объяснять услуги", "Делать сложные темы простыми: когда нужно УЗИ, анализы, вакцинация, чек-ап, стоматология.")
    add_card(slide, Inches(8.8), Inches(1.2), Inches(3.8), Inches(2.1), "3. Приводить в запись", "Каждый контент-юнит должен мягко вести в сообщения, звонок или запись на прием.")
    add_card(slide, Inches(0.7), Inches(3.65), Inches(5.9), Inches(2.5), "Что важно на старте", "Не пытаться делать все сразу. Нам нужен устойчивый ритм, понятные рубрики и база из фото, врачей и отзывов.")
    add_photo(slide, img_dir / "vet_exam.jpg", Inches(6.9), Inches(3.65), Inches(5.7), Inches(2.5))

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, LIGHT)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.5), "Роли площадок", 24, PRIMARY, True)
    add_card(slide, Inches(0.7), Inches(1.25), Inches(4.0), Inches(4.8), "Telegram", "Канал доверия и оперативных сообщений.\n\nПодходит для:\n- экспертных и сервисных постов\n- кейсов и новостей\n- прямых переходов в сообщения", fill="FFFFFF")
    add_card(slide, Inches(4.95), Inches(1.25), Inches(4.0), Inches(4.8), "VK", "Площадка для охвата, отзывов и прогрева.\n\nПодходит для:\n- постов и клипов\n- опросов, комментариев\n- входящих сообщений и акций", fill="FFFFFF")
    add_card(slide, Inches(9.2), Inches(1.25), Inches(3.4), Inches(4.8), "Instagram", "Визуальная витрина.\n\nПодходит для:\n- reels и каруселей\n- лиц врачей\n- stories и сохранений", fill="FFFFFF")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, "FFFFFF")
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.5), "Принципы ведения", 24, PRIMARY, True)
    y = 1.2
    for index, item in enumerate(PRINCIPLES, start=1):
        add_card(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(0.78), f"{index}.", item, fill="F9FAFB")
        y += 0.9

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, LIGHT)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.5), "Оптимальные рубрики", 24, PRIMARY, True)
    positions = [
        (Inches(0.7), Inches(1.3)),
        (Inches(6.6), Inches(1.3)),
        (Inches(0.7), Inches(4.0)),
        (Inches(6.6), Inches(4.0)),
    ]
    for (left, top), text in zip(positions, PILLARS):
        add_card(slide, left, top, Inches(5.7), Inches(2.2), text.split(":")[0], text.split(": ", 1)[1], fill="FFFFFF")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, "FFFFFF")
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.5), "Оптимальный ритм публикаций", 24, PRIMARY, True)
    add_photo(slide, img_dir / "ultrasound.jpg", Inches(8.0), Inches(1.3), Inches(4.6), Inches(5.5))
    y = 1.35
    for platform, frequency, desc in OPTIMAL_FORMATS:
        add_card(slide, Inches(0.8), Inches(y), Inches(6.7), Inches(1.45), f"{platform} - {frequency}", desc, fill="F8FAFC")
        y += 1.65
    add_textbox(slide, Inches(0.8), Inches(6.35), Inches(6.7), Inches(0.5), "Главное правило: один смысл перерабатываем под 3 площадки, а не делаем 3 разных контент-плана.", 11, MUTED)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, LIGHT)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.5), "Контент-план: неделя 1", 24, PRIMARY, True)
    week1 = CONTENT_PLAN[:7]
    y = 1.0
    for item in week1:
        body = f"{item['type']} • {item['formats']}\n{item['topic']}\nЦель: {item['goal']}"
        add_card(slide, Inches(0.7), Inches(y), Inches(12.0), Inches(0.68), f"{item['date'].strftime('%d.%m')} • {weekday_ru(item['date'])}", body, fill="FFFFFF")
        y += 0.8

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, LIGHT)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.5), "Контент-план: неделя 2", 24, PRIMARY, True)
    week2 = CONTENT_PLAN[7:]
    y = 1.0
    for item in week2:
        body = f"{item['type']} • {item['formats']}\n{item['topic']}\nЦель: {item['goal']}"
        add_card(slide, Inches(0.7), Inches(y), Inches(12.0), Inches(0.68), f"{item['date'].strftime('%d.%m')} • {weekday_ru(item['date'])}", body, fill="FFFFFF")
        y += 0.8

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, "FFFFFF")
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.5), "Что нужно для стабильной работы", 24, PRIMARY, True)
    y = 1.2
    for item in REQUIRED_ASSETS:
        add_card(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(0.72), "Нужно", item, fill="F9FAFB")
        y += 0.82

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_slide_bg(slide, PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5), "KPI первых 14 дней и следующий шаг", 24, "FFFFFF", True)
    y = 1.5
    for item in KPI:
        add_card(slide, Inches(0.8), Inches(y), Inches(7.1), Inches(0.82), "Метрика", item, fill="F3F4F6")
        y += 0.95
    add_photo(slide, img_dir / "dog_happy.jpg", Inches(8.4), Inches(1.5), Inches(4.1), Inches(3.3))
    add_textbox(slide, Inches(8.4), Inches(5.1), Inches(4.3), Inches(1.2), "Следующий шаг:\nутвердить рубрики, собрать фотоархив и снять 1 пакет контента на 2 недели вперед.", 15, ACCENT_SOFT, True)

    output = OUT_DIR / "Презентация_для_руководства_SMM.pptx"
    prs.save(output)
    return output


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    generated = [
        build_content_plan_docx(),
        build_templates_docx(),
        build_fill_forms_md(),
        build_presentation(),
    ]
    for path in generated:
        print(path.relative_to(ROOT))


if __name__ == "__main__":
    main()
